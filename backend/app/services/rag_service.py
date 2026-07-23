import os
import re
import time
import fitz
import numpy as np
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List, Dict, Any, Optional
from sqlalchemy.orm import Session
from sentence_transformers import SentenceTransformer
from app.core.config import settings
from app.database.connection import SessionLocal
from app.models.models import DocumentChunk

SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".txt", ".pptx")


class SimpleVectorStore:
    def __init__(self):
        self.model = SentenceTransformer("all-MiniLM-L6-v2")
        self.cache = {}

    def get_embeddings(self, texts: List[str]) -> np.ndarray:
        if not texts:
            return np.array([], dtype=np.float32)
        return self.model.encode(texts, convert_to_numpy=True)

    def search(self, query: str, chunks: List[DocumentChunk], top_k: int = 5, db: Optional[Session] = None) -> List[Dict[str, Any]]:
        if not chunks:
            return []

        query_vector = self.model.encode([query], convert_to_numpy=True)[0]
        chunk_embeddings = []
        chunks_to_update = []

        for chunk in chunks:
            if chunk.embedding is not None:
                emb = np.frombuffer(chunk.embedding, dtype=np.float32)
            else:
                emb = self.model.encode([chunk.content], convert_to_numpy=True)[0]
                chunk.embedding = emb.astype(np.float32).tobytes()
                chunks_to_update.append(chunk)
            chunk_embeddings.append(emb)

        if chunks_to_update and db:
            try:
                db.commit()
            except Exception as e:
                print(f"Error caching fallback embeddings: {e}")

        scores = []
        for i, emb in enumerate(chunk_embeddings):
            norm_q = np.linalg.norm(query_vector)
            norm_e = np.linalg.norm(emb)
            if norm_q == 0 or norm_e == 0:
                score = 0.0
            else:
                score = float(np.dot(query_vector, emb) / (norm_q * norm_e))
            scores.append((score, chunks[i]))

        scores.sort(key=lambda x: x[0], reverse=True)
        results = []
        seen = set()
        for score, chunk in scores[: top_k * 2]:
            content = re.sub(r"\s+", " ", chunk.content).strip()
            if not content or content in seen:
                continue
            seen.add(content)
            results.append({
                "score": score,
                "content": content,
                "page_number": chunk.page_number,
                "chunk_index": chunk.chunk_index,
            })
            if len(results) >= top_k:
                break
        return results


vector_store = SimpleVectorStore()


def extract_text_from_pdf(file_path: str) -> List[Dict[str, Any]]:
    start = time.perf_counter()
    doc = fitz.open(file_path)
    pages_content = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text()
        if text.strip():
            pages_content.append({
                "page_number": page_num + 1,
                "text": text,
            })
    elapsed_ms = round((time.perf_counter() - start) * 1000, 1)
    print(f"PDF Extraction: {elapsed_ms}ms")
    return pages_content

def extract_text_from_docx(file_path: str, words_per_page: int = 500) -> List[Dict[str, Any]]:
    start = time.perf_counter()
    doc = DocxDocument(file_path)
    pages_content = []
    current_words = []
    page_num = 1

    def flush():
        nonlocal current_words, page_num
        if current_words:
            pages_content.append({
                "page_number": page_num,
                "text": " ".join(current_words),
            })
            page_num += 1
            current_words = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        current_words.extend(text.split())
        if len(current_words) >= words_per_page:
            flush()

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                current_words.extend(row_text.split())
                if len(current_words) >= words_per_page:
                    flush()

    flush()

    elapsed_ms = round((time.perf_counter() - start) * 1000, 1)
    print(f"DOCX Extraction: {elapsed_ms}ms")
    return pages_content


def extract_text_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    start = time.perf_counter()
    prs = Presentation(file_path)
    pages_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        slide_text_parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_text_parts.append(row_text)

        text = "\n".join(slide_text_parts)
        if text.strip():
            pages_content.append({"page_number": slide_num, "text": text})

    elapsed_ms = round((time.perf_counter() - start) * 1000, 1)
    print(f"PPTX Extraction: {elapsed_ms}ms")
    return pages_content


def extract_text_from_txt(file_path: str, words_per_page: int = 500) -> List[Dict[str, Any]]:
    start = time.perf_counter()
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw_text = f.read()

    words = raw_text.split()
    pages_content = []
    page_num = 1
    for i in range(0, len(words), words_per_page):
        chunk_words = words[i:i + words_per_page]
        text = " ".join(chunk_words)
        if text.strip():
            pages_content.append({"page_number": page_num, "text": text})
            page_num += 1

    elapsed_ms = round((time.perf_counter() - start) * 1000, 1)
    print(f"TXT Extraction: {elapsed_ms}ms")
    return pages_content


def extract_text(file_path: str) -> List[Dict[str, Any]]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    raise ValueError(f"Unsupported file type: {ext}")


def chunk_text(pages: List[Dict[str, Any]], chunk_size: int = 700, chunk_overlap: int = 120) -> List[Dict[str, Any]]:
    chunks = []
    chunk_index = 0
    for page in pages:
        text = page["text"]
        page_num = page["page_number"]

        words = text.split()
        current_chunk_words = []
        current_length = 0

        for word in words:
            current_chunk_words.append(word)
            current_length += len(word) + 1
            if current_length >= chunk_size:
                chunk_text = " ".join(current_chunk_words)
                chunks.append({
                    "content": chunk_text,
                    "page_number": page_num,
                    "chunk_index": chunk_index,
                })
                chunk_index += 1

                overlap_words = current_chunk_words[-max(1, len(current_chunk_words) // 5):]
                current_chunk_words = list(overlap_words)
                current_length = sum(len(w) + 1 for w in current_chunk_words)

        if current_chunk_words:
            chunk_text = " ".join(current_chunk_words)
            chunks.append({
                "content": chunk_text,
                "page_number": page_num,
                "chunk_index": chunk_index,
            })
            chunk_index += 1

    return chunks


def process_and_index_document(file_path: str, document_id: int) -> None:
    start = time.perf_counter()
    pages = extract_text(file_path)
    chunks_data = chunk_text(pages)

    db = SessionLocal()
    try:
        texts = [c["content"] for c in chunks_data]
        embeddings = vector_store.get_embeddings(texts) if texts else []

        for i, c in enumerate(chunks_data):
            emb_bytes = embeddings[i].astype(np.float32).tobytes() if i < len(embeddings) else None
            chunk_model = DocumentChunk(
                document_id=document_id,
                content=c["content"],
                page_number=c["page_number"],
                chunk_index=c["chunk_index"],
                embedding=emb_bytes,
            )
            db.add(chunk_model)
        db.commit()
    finally:
        db.close()

    elapsed_ms = round((time.perf_counter() - start) * 1000, 1)
    print(f"Document Indexing Total: {elapsed_ms}ms")
